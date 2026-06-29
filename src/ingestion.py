import os
import base64
import fitz  # PyMuPDF
import mimetypes
from docx import Document
from pptx import Presentation
from src import config
from src import llm

class RecursiveCharacterTextSplitter:
    def __init__(self, chunk_size=1000, chunk_overlap=200, separators=None):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.separators = separators or ["\n\n", "\n", " ", ""]

    def split_text(self, text: str) -> list[str]:
        """Splits a single string into chunks recursively using configured separators."""
        if not text:
            return []
            
        def split_rec(txt, separators):
            if len(txt) <= self.chunk_size:
                return [txt]
            
            if not separators:
                # Hard split by character counts if no separators left
                step = max(1, self.chunk_size - self.chunk_overlap)
                return [txt[i:i+self.chunk_size] for i in range(0, len(txt), step)]
                
            separator = separators[0]
            splits = txt.split(separator)
            
            chunks = []
            current_chunk = []
            current_len = 0
            
            for part in splits:
                # If a single split part is still larger than the chunk size, split it with remaining separators
                sub_parts = []
                if len(part) > self.chunk_size:
                    sub_parts = split_rec(part, separators[1:])
                else:
                    sub_parts = [part]
                    
                for sub_part in sub_parts:
                    sep_len = len(separator) if current_chunk else 0
                    if current_len + len(sub_part) + sep_len <= self.chunk_size:
                        current_chunk.append(sub_part)
                        current_len += len(sub_part) + sep_len
                    else:
                        if current_chunk:
                            chunks.append(separator.join(current_chunk))
                        
                        # Apply overlap by tracing back
                        overlap_chunk = []
                        overlap_len = 0
                        for c in reversed(current_chunk):
                            s_len = len(separator) if overlap_chunk else 0
                            if overlap_len + len(c) + s_len <= self.chunk_overlap:
                                overlap_chunk.insert(0, c)
                                overlap_len += len(c) + s_len
                            else:
                                break
                        
                        current_chunk = overlap_chunk + [sub_part]
                        current_len = sum(len(x) for x in current_chunk) + len(separator) * (len(current_chunk) - 1)
            
            if current_chunk:
                chunks.append(separator.join(current_chunk))
            return chunks

        return split_rec(text, self.separators)

def parse_pdf(file_path: str) -> list[dict]:
    """
    Parses a PDF file using PyMuPDF, extracts its text, and returns a list of chunks:
    [{"text": chunk_text, "source": filename, "page": page_number}, ...]
    """
    filename = os.path.basename(file_path)
    chunks = []
    
    try:
        doc = fitz.open(file_path)
        total_pages = len(doc)
        
        # Check if the PDF has actual extractable text
        is_scanned = True
        full_extracted_text = ""
        
        pages_text = []
        for page in doc:
            text = page.get_text() or ""
            pages_text.append(text)
            full_extracted_text += text
            
        if len(full_extracted_text.strip()) > 100:
            is_scanned = False
            
        if is_scanned:
            # Document is scanned. Try using model's multimodal capability (Gemini API)
            active_cfg = llm.get_active_config()
            if active_cfg["provider"] == "gemini" and active_cfg["gemini_key"]:
                # Use Gemini multimodal API to read and transcribe PDF pages
                from google import genai
                from google.genai import types
                
                client = genai.Client(api_key=active_cfg["gemini_key"])
                with open(file_path, "rb") as f:
                    pdf_bytes = f.read()
                
                # Ask Gemini to transcribe the document
                prompt = "Extract and transcribe all the textual content from this scanned document. Maintain the layout, paragraphs, and tables as closely as possible."
                response = client.models.generate_content(
                    model=llm.GEMINI_TASK_MODELS.get("vision", "gemini-2.5-flash"),
                    contents=[
                        prompt,
                        types.Part.from_bytes(data=pdf_bytes, mime_type="application/pdf")
                    ]
                )
                
                # Split the transcribed text
                splitter = RecursiveCharacterTextSplitter(
                    chunk_size=config.CHUNK_SIZE,
                    chunk_overlap=config.CHUNK_OVERLAP
                )
                raw_chunks = splitter.split_text(response.text)
                
                for idx, chunk in enumerate(raw_chunks):
                    chunks.append({
                        "text": chunk,
                        "source": filename,
                        "page": 1  # OCRed output is returned as a single response
                    })
            else:
                # Local Ollama or no keys. Fallback: raise an exception with instructions
                raise ValueError(
                    f"The PDF '{filename}' appears to be a scanned document (no digital text found). "
                    "To ingest scanned documents, please configure a Gemini API Key in the sidebar to enable automated vision-based OCR transcription, or upload a digital PDF."
                )
        else:
            # Digital PDF. Process page-by-page
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=config.CHUNK_SIZE,
                chunk_overlap=config.CHUNK_OVERLAP
            )
            
            for i, page_text in enumerate(pages_text):
                page_num = i + 1
                if not page_text.strip():
                    continue
                
                raw_chunks = splitter.split_text(page_text)
                for chunk in raw_chunks:
                    chunks.append({
                        "text": chunk,
                        "source": filename,
                        "page": page_num
                    })
                    
    except Exception as e:
        print(f"Error parsing PDF {file_path}: {e}")
        raise e
        
    return chunks

def parse_docx(file_path: str) -> list[dict]:
    """
    Parses a DOCX file, extracts paragraphs and table cells, and returns a list of chunks:
    [{"text": chunk_text, "source": filename, "page": 1}, ...]
    """
    filename = os.path.basename(file_path)
    try:
        doc = Document(file_path)
        full_text = []
        
        # Extract paragraph text
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text.strip())
                
        # Extract table cells
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    full_text.append(" | ".join(row_text))
                    
        joined_text = "\n\n".join(full_text)
        
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=config.CHUNK_SIZE,
            chunk_overlap=config.CHUNK_OVERLAP
        )
        raw_chunks = splitter.split_text(joined_text)
        
        chunks = []
        for idx, chunk in enumerate(raw_chunks):
            chunks.append({
                "text": chunk,
                "source": filename,
                "page": 1  # DOCX doesn't have standard physical pages in python-docx
            })
        return chunks
    except Exception as e:
        print(f"Error parsing DOCX {file_path}: {e}")
        raise e

def parse_pptx(file_path: str) -> list[dict]:
    """
    Parses a PPTX file, extracts shape text slide-by-slide, and returns a list of chunks:
    [{"text": chunk_text, "source": filename, "page": slide_number}, ...]
    """
    filename = os.path.basename(file_path)
    try:
        prs = Presentation(file_path)
        chunks = []
        
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=config.CHUNK_SIZE,
            chunk_overlap=config.CHUNK_OVERLAP
        )
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1
            slide_text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_parts.append(shape.text.strip())
                    
            slide_text = "\n".join(slide_text_parts)
            if not slide_text.strip():
                continue
                
            raw_chunks = splitter.split_text(slide_text)
            for chunk in raw_chunks:
                chunks.append({
                    "text": chunk,
                    "source": filename,
                    "page": slide_num
                })
        return chunks
    except Exception as e:
        print(f"Error parsing PPTX {file_path}: {e}")
        raise e

def parse_image(file_path: str) -> list[dict]:
    """
    Uses the configured provider's vision model to describe an image and returns a single chunk.
    """
    filename = os.path.basename(file_path)
    try:
        with open(file_path, "rb") as f:
            image_bytes = f.read()
            
        mime_type, _ = mimetypes.guess_type(file_path)
        if not mime_type:
            mime_type = "image/jpeg"
            
        prompt = (
            "Analyze this image in detail. Generate a comprehensive textual description. "
            "Describe all visual entities, objects, texts, layouts, processes, and relationships visible in the image. "
            "Be precise, detailed, and factual. Your description will be parsed for knowledge engineering."
        )
        
        description = llm.generate_vision(image_bytes, mime_type, prompt)
        
        return [{
            "text": description,
            "source": filename,
            "page": 1
        }]
    except Exception as e:
        print(f"Error parsing image {file_path}: {e}")
        raise e

def ingest_file(file_path: str) -> list[dict]:
    """
    Ingests any supported file type (.pdf, .docx, .pptx, .jpg, .jpeg, .png)
    and returns a uniform list of chunks:
    [{"text": chunk_text, "source": filename, "page": page_number}, ...]
    """
    _, ext = os.path.splitext(file_path.lower())
    if ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    elif ext in [".jpg", ".jpeg", ".png"]:
        return parse_image(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")
